import os
import io
# docx
from docx import Document
from docx.shared import Inches

# pptx
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches
# docx2html or docx2md
from mammoth import convert as mammoth_convert

from ..cat import PDF4Cat

class any_doc_convert(PDF4Cat):
	
	"""Subclass of PDF4Cat parent class
	
	Args:
		doc_file (None, optional): Document file (for multiple operations, 'use input_doc_list')
		input_doc_list (list, optional): List of input docs
		passwd (str, optional): Document password (for crypt/decrypt)
		progress_callback (None, optional): Progress callback like:
	
	Raises:
		TypeError: If you use doc_file with input_doc_list (you can use only one)
	"""
	
	def __init__(self, *args, **kwargs):
		super(any_doc_convert, self).__init__(*args, **kwargs)

	@PDF4Cat.run_in_subprocess
	# need add html to pdf
	def convert2pdf(self, output_pdf, use_soffice=False): # only for text based pdf and works not for all
		"""Pdf to any (using PyMuPDF or Libre Office)
		
		Args:
			output_pdf (None, optional): Output pdf file
			use_soffice (bool, optional): Use Libre Office converter
		"""
		if isinstance(output_pdf, str):
			ext = os.path.splitext(self.doc_filename)[1][1:]
			if ext in ['docx', 'doc'] and not use_soffice:
				self.docx2pdf(output_pdf)
				return
			elif ext in ['png', 'jpg'] and not use_soffice:
				self.img2pdf(output_pdf)
			elif "."+ext in self.libre_exts and use_soffice:
				self.soffice_convert2pdf(output_pdf)
				return
			else:
				raise NotImplementedError(f"File extension '{self.doc_fileext}' => '.pdf' not supported")
			# elif ext == 'pptx' or ext == 'ppt':
			# 	self.pptx2pdf(output_pdf)
			# 	return
		
		doc = self.pdf_open(self.doc_file)

		b = doc.convert_to_pdf()  # convert to pdf
		pdf = self.pdf_open("pdf", b)  # open as pdf

		# ***
		# toc = doc.het_toc()  # table of contents of input
		# pdf.set_toc(toc)  # simply set it for output
		# meta = doc.metadata  # read and set metadata
		# if not meta["producer"]:
		# 	meta["producer"] = "PDF4Cat https://github.com/BlackCatDevel0per/PDF4Cat"

		# if not meta["creator"]:
		# 	meta["creator"] = "PDF4Cat pdf tool"
		# meta["modDate"] = self.fitz_get_pdf_now()
		# meta["creationDate"] = meta["modDate"]
		# pdf.set_metadata(meta)
		# ***

		# now process the links
		link_cnti = 0
		link_skip = 0
		for pinput in doc:  # iterate through input pages
			links = pinput.get_links()  # get list of links
			link_cnti += len(links)  # count how many
			pout = pdf[pinput.number]  # read corresp. output page
			for l in links:  # iterate though the links
				if l["kind"] == self.fitz_LINK_NAMED:  # we do not handle named links
					print("named link page", pinput.number, l)
					link_skip += 1  # count them
					continue
				pout.insert_link(l)  # simply output the others

		# save the conversion result
		pdf.save(output_pdf, garbage=4, deflate=True)
		# say how many named links we skipped
		if link_cnti > 0:
			print("Skipped %i named links of a total of %i in input." % (link_skip, link_cnti))
	
	# Generate name with BytesIO object (it is faster)
	def gen_images4conv(self, pdf) -> bytes:
		"""Generator, generate BytesIO object
		
		Args:
			pdf (None, optional): pdf object (PDF4Cat.open)
		
		Yields:
			bytes: BytesIO
		"""
		noOfPages = range(pdf.page_count)

		for pageNo in noOfPages:
			io_data = io.BytesIO()
			#
			page = pdf.load_page(pageNo) #number of page
			pix = page.get_pixmap()
			del page
			io_data.write(pix.tobytes(output="png"))
			del pix
			#

			imfi = io_data
			yield imfi

	@PDF4Cat.run_in_subprocess
	def pdf2pptx(self, output_pptx, A4=True):
		"""Pdf to pptx (using PyMuPDF)
		
		Args:
			output_pptx (None, optional): Output pptx file
			A4 (bool, optional): Use Inches for A4 page
		"""
		if not output_pptx:
			output_pptx = os.path.join(self.doc_path, self.doc_name+"_out.pdf")
		output_pptx = os.path.join(os.getcwd(), output_pptx)

		pdf = self.pdf_open(self.doc_file, passwd=self.passwd)
		prs = Presentation()
		w, h = 13.333, 7.5
		if A4:
			prs.slide_height=Inches(11)
			prs.slide_width=Inches(8.5)
			w, h = 8.5, 11

		blank_slide_layout = prs.slide_layouts[6]

		for io_data in self.gen_images4conv(pdf):
			slide = prs.slides.add_slide(blank_slide_layout)
			# slide.shapes.add_picture(io_data, 0, 0, width=Inches(13.333), height=Inches(7.5))
			slide.shapes.add_picture(io_data, 0, 0, width=Inches(w), height=Inches(h))
			del io_data
			self.counter += 1 #need enumerate
			self.progress_callback(self.counter, pdf.page_count)

		prs.save(output_pptx)
		self.counter = 0

	@PDF4Cat.run_in_subprocess
	def pdf2docx(self, output_docx):
		"""Pdf to docx (using PyMuPDF)
		
		Args:
			output_docx (None, optional): Output docx file
		"""
		if not output_docx:
			output_docx = os.path.join(self.doc_path, self.doc_name+"_out.pdf")
		output_docx = os.path.join(os.getcwd(), output_docx)

		pdf = self.pdf_open(self.doc_file, passwd=self.passwd)
		document = Document()

		for io_data in self.gen_images4conv(pdf):
			# document.add_picture(io_data, width=Inches(8.5), height=Inches(11))
			document.add_picture(io_data, width=Inches(5.7), height=Inches(9)) # 5.75
			del io_data
			self.counter += 1 #need enumerate
			self.progress_callback(self.counter, pdf.page_count)

		document.save(output_docx)
		self.counter = 0

	@PDF4Cat.run_in_subprocess
	def docx2html(self, output_doc, style_map = None):
		"""docx to html (using PyMuPDF)
		
		Args:
			output_html (None, optional): Output html file
		"""
		if not output_doc:
			output_doc = os.path.join(self.doc_path, self.doc_name+"_out.pdf")
		if isinstance(output_doc, str):
			output_doc = os.path.join(os.getcwd(), output_doc)

		if not style_map:
			style_map = None
		else:
			with open(style_map) as style_map_fileobj:
				style_map = style_map_fileobj.read()

		with open(self.doc_file, "rb") as docx_file:
			result = mammoth_convert(
			docx_file,
			style_map=style_map,
			output_format='html', # markdown have too
		)
		# ***
		if not isinstance(output_doc, str): # need os like path..
			output_doc.write(result.value.encode())
			return
		# ***
		with open(output_doc, "wb") as output:
			output.write(result.value.encode())


	@PDF4Cat.run_in_subprocess
	def docx2pdf(self, output_pdf):
		"""docx to pdf (using PyMuPDF [docx=>html=>pdf])
		
		Args:
			output_pdf (None, optional): Output pdf file
		"""
		if not output_pdf:
			output_pdf = os.path.join(self.doc_path, self.doc_name+"_out.pdf")
		output_pdf = os.path.join(os.getcwd(), output_pdf)
		html_tmp = io.BytesIO()

		self.docx2html(html_tmp, run_in_subprocess=False)
		doc = self.pdf_open(filename="html", stream=html_tmp)
		pdfbytes = doc.convert_to_pdf() # need convert to image..?
		doc.close()
		with open(output_pdf, 'wb') as output:
			output.write(pdfbytes)
		# pdf = self.pdf_open("pdf", pdfbytes)
		# pdf.save(output_pdf)
		# pdf.close()

			